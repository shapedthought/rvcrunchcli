# RvCruncher CLI

import argparse
import pandas as pd
import numpy as np
import math
import json
import pprint
from workload_dict import work_dict
from pptx import Presentation


def main():
    """
    RvCrunch- simple reporter

    RvCrunch is a simple tool to provide a report that has relevant information for a
    Veeam sizing. The base function produces a report with an environment breakdown.
    There are also mini reports that provide quick reports on key information from
    the RvTools report.

    The tool accepts RvTools files in excel format (.xlsx)

    The script requries 'pandas' and 'numpy' to be installed within the Python environment
    for the script to run.

    Specifying the file path as the first arguement will run the main "Crunch" report.
    Using the optional parameter -sr or --skipreport will skip this report.

    Using -po or --poweredon will only display Powered On virtual machine info.

    Using -vse or --vseexport will export a VSE compatible txt file that can be uploaded.
    Note that the will be Powered On VMs only.

    Using -sn or --savename allows you to specify the report name.

    There are also several mini reports that provide certain key info. Remember to use the -sr flag
    if you do not want to main report to run.
    -vm : Provides the total VM count
    -dc : Provides the Datacenter names
    -cl : Provides the Cluster names
    -dc_cl : Groups the clusters by the Datacenter
    -rdm : Export an excel with RDMs by VM plus capacity in TB
    Note: RDM report is allocated capacity, mapping vDisk to vPartition isn't possible

    """
    parser = argparse.ArgumentParser(
        prog="RvCrunch",
        usage="%(prog)s [OPTIONS] path",
        description="Provides relevant info from RvTools files",
    )
    parser.add_argument("Path", metavar="path", type=str, help="the path to the file")
    mini_reports = parser.add_argument_group("Mini Reports")
    parser.add_argument(
        "-sr",
        "--skipreport",
        action="store_false",
        help="Skips the main report generation",
    )
    parser.add_argument(
        "-po", "--poweredon", action="store_true", help="Just Powered On VMs in report"
    )
    parser.add_argument("-sn", "--savename", type=str, help="Report name")
    parser.add_argument(
        "-vse",
        "--vseexport",
        action="store_true",
        help="Creates VSE compatible txt file",
    )
    mini_reports.add_argument(
        "-vm",
        "--vmcount",
        action="store_true",
        help="Total quantity of VMs, powered on and off, by datacenter and cluster",
    )
    mini_reports.add_argument(
        "-dc", "--datacenter", action="store_true", help="Show datacenters"
    )
    mini_reports.add_argument(
        "-cl", "--clusters", action="store_true", help="Show clusters"
    )
    mini_reports.add_argument(
        "-dc_cl", "--dc_clusters", action="store_true", help="Show clusters per-DC"
    )
    mini_reports.add_argument(
        "-rdm", "--rdmcheck", action="store_true", help="Check for RDMs"
    )

    args = parser.parse_args()

    # Importing the data from the RvTools
    try:
        print("Running... please wait")
        v_info_data = pd.read_excel(args.Path, sheet_name="vInfo")
        v_partition_data = pd.read_excel(args.Path, sheet_name="vPartition")
        v_info_data["VM Count"] = 1
    except:
        raise FileNotFoundError("File or Required tabs not present")

    # Main report code
    if args.skipreport:
        try:
            # Filter out all Powered Off VMs
            if args.poweredon:
                v_info_data = v_info_data[v_info_data["Powerstate"] == "poweredOn"]
                v_partition_data = v_partition_data[
                    v_partition_data["Powerstate"] == "poweredOn"
                ]

            # Add TB columns
            v_info_data["vInfo TB"] = round(v_info_data["In Use MB"] / (1024 ** 2), 2)
            v_partition_data["vPartition TB"] = round(
                (v_partition_data["Capacity MB"] - v_partition_data["Free MB"])
                / (1024 ** 2),
                2,
            )
            v_partition_data = v_partition_data[["VM", "vPartition TB"]]

            # Use groupby to calculate the total vpartition capacity per-vm
            v_partition_data = v_partition_data.groupby(["VM"], as_index=False).sum()

            # Combine the data into a single dataframe
            combo_data = pd.merge(
                v_info_data, v_partition_data, on="VM", how="outer"
            ).fillna(0)

            # Replace the vpartition capacity with the vinfo capacity where the vpartiion capacity is zero
            combo_data["Agg Cap TB"] = np.where(
                combo_data["vPartition TB"] == 0,
                combo_data["vInfo TB"],
                combo_data["vPartition TB"],
            )
            combo_data = combo_data[
                [
                    "Powerstate",
                    "VM Count",
                    "Disks",
                    "vInfo TB",
                    "vPartition TB",
                    "Datacenter",
                    "Agg Cap TB",
                    "Cluster",
                ]
            ]

            # Save the data to a file
            save_name = (
                "report.xlsx" if args.savename == None else args.savename + ".xlsx"
            )
            print(f"Crunch file saved as {save_name}")
            cap_data = combo_data.groupby(["Powerstate", "Datacenter", "Cluster"]).sum()
            cap_data.to_excel(save_name)
        except:
            raise RuntimeError("Error in calculation")
    # Presentation test
    # pr1 = Presentation()
    # slide1_register = pr1.slide_layouts[0]
    # slide1 = pr1.slides.add_slide(slide1_register)
    # title1 = slide1.shapes.title
    # title1.text = 'Veeam Report'
    # subtitle1 = slide1.placeholders[1]
    # subtitle1.text = 'RvTools Info'
    # pr1.save('Veeam_report.pptx')

    if args.vmcount:
        # VM count mini report
        vm_count = v_info_data.groupby(["Powerstate"])["VM Count"].apply(sum)
        print("")
        print("***VM Count***")
        print(vm_count)
    if args.datacenter:
        # Datacenter mini report
        datacenters = pd.unique(v_info_data["Datacenter"])
        print("")
        print("***Datacenters***")
        for datacenter in datacenters:
            print(datacenter)
    if args.clusters:
        # Cluster mini report
        clusters = pd.unique(v_info_data["Cluster"])
        print("")
        print("***Clusters***")
        for cluster in clusters:
            print(cluster)
    if args.dc_clusters:
        # Datacetre and cluster mini report
        v_info_data["vInfo TB"] = round(v_info_data["In Use MB"] / (1024 ** 2), 2)
        dc_cluster = v_info_data.groupby(["Powerstate", "Datacenter", "Cluster"])[
            ["VM Count", "vInfo TB"]
        ].apply(sum)
        print("")
        print("***DC and Clusters***")
        print(dc_cluster)
    if args.rdmcheck:
        # RDM check mini report
        v_disk_data = pd.read_excel(args.Path, sheet_name="vDisk").fillna("None")
        v_disk_data["Capacity TB"] = round(v_disk_data["Capacity MB"] / (1024 ** 2), 2)
        result = (
            v_disk_data[v_disk_data["Raw Comp. Mode"].str.contains("physicalMode")]
            .groupby(["VM", "Raw Comp. Mode"], as_index=False)
            .agg({"Disk": "count", "Capacity TB": "sum"})
        )
        if len(result) > 0:
            print("")
            print("***RDM Report***")
            print(result)
            confirm = input("Export to Excel? Y/n: ")
            if confirm == "Y":
                result.to_excel("RDM_report.xlsx")
                print("RDM Report saved")
        else:
            print("No RDMs in report")

    def object_creator(workloads):
        """
        Helper function that creates the modified version of the VSE workloads object
        """
        new_dict = work_dict.copy()
        new_dict["workLoadName"] = workloads["workLoadName"]
        new_dict["workLoadCap"] = round(workloads["workLoadCap"], 2)
        new_dict["vmQty"] = workloads["vmQty"]
        new_dict["vmdkQty"] = workloads["vmdkQty"]
        new_dict["vmVmdkRatio"] = math.ceil(workloads["vmdkQty"] / workloads["vmQty"])
        return new_dict

    def vse_export():
        """
        Re-maps the RvTools dataframe to the correct VSE format.
        """
        results = []
        # filering out the powered off, the group by datacentre
        combo_powered_on = (
            combo_data[combo_data["Powerstate"] == "poweredOn"]
            .groupby(["Datacenter"], as_index=False)
            .sum()
        )
        # Loop through the dataframe, create a new dict to use in the
        # object_creator function.
        for _, row in combo_powered_on.iterrows():
            workload = {
                "workLoadName": row["Datacenter"],
                "workLoadCap": row["Agg Cap TB"],
                "vmQty": row["VM Count"],
                "vmdkQty": row["Disks"],
            }
            new_object = object_creator(workload)
            results.append(new_object)
        export_dict = dict({"workload": results})
        json_dict = json.dumps(export_dict)
        f = open("vse_export.txt", "w")
        f.write(json_dict)
        f.close
        print("VSE File Exported Successfully")

    if args.vseexport:
        try:
            vse_export()
        except:
            raise RuntimeError("Something bad happend...")


if __name__ == "__main__":
    main()